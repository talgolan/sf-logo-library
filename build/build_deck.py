"""Build a PPTX showcasing Salesforce, MuleSoft, and Informatica logos from dam.usefulto.me.

Proportions are preserved by setting only pic.width; python-pptx derives pic.height
from the image's intrinsic aspect ratio.
"""
import json
import subprocess
import urllib.parse
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

BASE = "https://dam.usefulto.me/"
BUILD = Path(__file__).parent
LOGO_DIR = BUILD / "logos"
LOGO_DIR.mkdir(exist_ok=True)

MANIFEST_URL = BASE + "manifest.json"
TARGET_BRANDS = {"salesforce", "mulesoft", "informatica"}


def fetch_manifest() -> dict:
    with urllib.request.urlopen(MANIFEST_URL) as r:
        return json.load(r)


def download(rel_path: str) -> Path:
    url = BASE + urllib.parse.quote(rel_path)
    local = LOGO_DIR / rel_path.replace("/", "__")
    if not local.exists():
        print(f"  fetch {rel_path}")
        with urllib.request.urlopen(url) as r, open(local, "wb") as f:
            f.write(r.read())
    return local


def svg_to_png(svg_path: Path, width_px: int = 2000) -> Path:
    png_path = svg_path.with_suffix(".converted.png")
    if not png_path.exists():
        print(f"  convert {svg_path.name} -> PNG @ {width_px}px")
        subprocess.run(
            [
                "magick",
                "-background", "none",
                "-density", "600",
                str(svg_path),
                "-resize", f"{width_px}x",
                str(png_path),
            ],
            check=True,
        )
    return png_path


def ensure_png(logo: dict) -> Path:
    if logo.get("png"):
        return download(logo["png"])
    svg = download(logo["svg"])
    return svg_to_png(svg)


# ---- build ---------------------------------------------------------------

manifest = fetch_manifest()
brands = [b for b in manifest["brands"] if b["id"] in TARGET_BRANDS]
# preserve requested ordering: salesforce, mulesoft, informatica
order = ["salesforce", "mulesoft", "informatica"]
brands.sort(key=lambda b: order.index(b["id"]))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

DARK_BG = RGBColor(0x03, 0x2D, 0x60)  # Salesforce deep navy-ish, works for all knockouts
LIGHT_BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x18, 0x1B, 0x1F)
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_textbox(slide, text: str, top: Emu, color: RGBColor, size_pt: int = 20, bold: bool = False) -> None:
    tb = slide.shapes.add_textbox(Inches(0.5), top, SW - Inches(1.0), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 2  # center
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_logo_centered(slide, png_path: Path, max_w_in: float, max_h_in: float) -> None:
    """Add picture centered in the slide, constrained to max_w_in x max_h_in,
    while preserving the image's intrinsic aspect ratio."""
    from PIL import Image
    with Image.open(png_path) as im:
        iw, ih = im.size
    img_ratio = iw / ih
    box_ratio = max_w_in / max_h_in
    if img_ratio >= box_ratio:
        # width-constrained
        w = Inches(max_w_in)
        pic = slide.shapes.add_picture(str(png_path), 0, 0, width=w)
    else:
        h = Inches(max_h_in)
        pic = slide.shapes.add_picture(str(png_path), 0, 0, height=h)
    # center
    pic.left = Emu(int((SW - pic.width) / 2))
    pic.top = Emu(int((SH - pic.height) / 2))


# ---- Title slide ---------------------------------------------------------
title = prs.slides.add_slide(BLANK)
set_slide_bg(title, LIGHT_BG)
add_textbox(title, "Salesforce · MuleSoft · Informatica", Inches(2.7), TEXT_DARK, size_pt=44, bold=True)
add_textbox(title, "Logo Reference — dam.usefulto.me", Inches(3.9), RGBColor(0x5E, 0x69, 0x75), size_pt=20)
add_textbox(title, f"Version {manifest.get('version', '')}", Inches(6.7), RGBColor(0x9A, 0xA2, 0xAE), size_pt=12)

# ---- Brand section + per-logo slides ------------------------------------
BRAND_DISPLAY = {"salesforce": "Salesforce", "mulesoft": "MuleSoft", "informatica": "Informatica"}

for brand in brands:
    # divider slide
    div = prs.slides.add_slide(BLANK)
    set_slide_bg(div, LIGHT_BG)
    add_textbox(div, BRAND_DISPLAY[brand["id"]], Inches(3.2), TEXT_DARK, size_pt=60, bold=True)
    add_textbox(div, f"{len(brand['logos'])} logo variants", Inches(4.3), RGBColor(0x5E, 0x69, 0x75), size_pt=18)

    for logo in brand["logos"]:
        png = ensure_png(logo)
        slide = prs.slides.add_slide(BLANK)
        is_dark = logo.get("background") == "dark"
        set_slide_bg(slide, DARK_BG if is_dark else LIGHT_BG)
        caption_color = TEXT_LIGHT if is_dark else TEXT_DARK
        subcaption_color = RGBColor(0xB8, 0xC7, 0xE0) if is_dark else RGBColor(0x5E, 0x69, 0x75)

        # logo area: center 70% of slide height, 80% width
        add_logo_centered(slide, png, max_w_in=10.0, max_h_in=4.8)

        # captions
        add_textbox(slide, logo["name"], Inches(0.55), caption_color, size_pt=24, bold=True)
        variant_line = f"{logo.get('variant', '')}  ·  {logo.get('type', '')}  ·  {logo.get('background', '')} background"
        add_textbox(slide, variant_line, Inches(1.15), subcaption_color, size_pt=12)

        ar = logo.get("aspect_ratio", {})
        dims = logo.get("dimensions", {})
        foot = f"{dims.get('width', '?')}×{dims.get('height', '?')} · ratio {ar.get('ratio', '?')} ({ar.get('decimal', '?')})"
        add_textbox(slide, foot, Inches(6.7), subcaption_color, size_pt=11)


out = BUILD.parent / "Salesforce-MuleSoft-Informatica-Logos.pptx"
prs.save(out)
print(f"\nwrote {out}  ({out.stat().st_size/1024:.1f} KB, {len(prs.slides)} slides)")
